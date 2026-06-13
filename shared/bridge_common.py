#!/usr/bin/env python3
"""bridge_common.py — MAGI の LLM ブリッジで共有するプロバイダ非依存ヘルパー群。

ask_openai.py(ChatGPT)と ask_gemini.py(Gemini)が共通で使う:
  - API キー解決(環境変数 → ルートの .env → .claude/settings.local.json の env、
    placeholder 除外)。シークレットは原則ルートの .env に集約する
  - 人格定義のフロントマター/先頭HTMLコメント除去(本文=システムプロンプト化)
  - 添付ファイルの種別判定とローカル抽出(Office)/読み込み(テキスト)
  - base64 化・MIME 推定・HTTP(JSON / multipart)ヘルパー

依存: 標準ライブラリのみ(Office 抽出時だけ python-docx / openpyxl / python-pptx)。
"""
import base64
import json
import mimetypes
import os
import random
import socket
import sys
import time
import urllib.error
import urllib.request

# Windows コンソール/パイプ経由でも日本語を UTF-8 で入出力できるよう固定する。
# stdin を含めるのが要点: パイプで渡したラウンド入力が OS 既定コーデック(例: cp932)で
# 復号されると不正なサロゲートが生じ、OpenAI 等が「unpaired UTF-16 surrogate」で 400 を返す。
for _stream in (sys.stdin, sys.stdout, sys.stderr):
    try:
        _stream.reconfigure(encoding="utf-8")
    except (AttributeError, ValueError):
        pass

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".webp"}
PDF_EXTS = {".pdf"}
OFFICE_EXTS = {".docx", ".xlsx", ".pptx"}
TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".log", ".yaml", ".yml"}


# --------------------------------------------------------------------------- #
# 設定(API キー)の解決
# --------------------------------------------------------------------------- #
def settings_env():
    """.claude/settings.local.json の env ブロックを読む(あれば)。
    Claude Code は起動時にしか settings の env を反映しないため、セッション中に
    キーを設定しても効くよう、フォールバックとして直接読み込む。"""
    here = os.path.dirname(os.path.abspath(__file__))
    path = os.path.join(here, os.pardir, ".claude", "settings.local.json")
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f).get("env", {}) or {}
    except (OSError, ValueError):
        return {}


def parse_env_file(path):
    """dotenv 形式(KEY=VALUE)のファイルを dict にする。標準ライブラリのみで実装。
    コメント行(#)・空行・export 接頭辞を許容し、値の両端の引用符は剥がす。"""
    env = {}
    try:
        with open(path, "r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line or line.startswith("#") or "=" not in line:
                    continue
                if line.startswith("export "):
                    line = line[len("export "):]
                key, _, val = line.partition("=")
                key, val = key.strip(), val.strip()
                if len(val) >= 2 and val[0] == val[-1] and val[0] in ("'", '"'):
                    val = val[1:-1]
                if key:
                    env[key] = val
    except OSError:
        return {}
    return env


def dotenv_env():
    """リポジトリルートの .env を読む(あれば)。シークレットの正本はここに集約する。"""
    here = os.path.dirname(os.path.abspath(__file__))
    return parse_env_file(os.path.join(here, os.pardir, ".env"))


def get_setting(*names):
    """環境変数 → ルートの .env → settings.local.json の env の順に解決する。
    複数の名前を渡すと先に見つかったものを返す(例: GEMINI_API_KEY, GOOGLE_API_KEY)。
    placeholder("...REPLACE...")・空文字は未設定扱い。"""
    for name in names:
        val = os.environ.get(name)
        if val:
            return val
    for source in (dotenv_env(), settings_env()):
        for name in names:
            val = source.get(name)
            if val and "REPLACE" not in val:
                return val
    return None


# --------------------------------------------------------------------------- #
# 人格定義 / 入力の読み込み
# --------------------------------------------------------------------------- #
def strip_frontmatter(text):
    """先頭の YAML フロントマター(--- ... ---)と、その後に続く先頭の HTML コメントを
    取り除き、人格本文(システムプロンプト)だけを返す。"""
    if text.startswith("---"):
        parts = text.split("\n")
        if parts and parts[0].strip() == "---":
            for i in range(1, len(parts)):
                if parts[i].strip() == "---":
                    text = "\n".join(parts[i + 1 :]).lstrip("\n")
                    break
    stripped = text.lstrip()
    while stripped.startswith("<!--"):
        end = stripped.find("-->")
        if end == -1:
            break
        stripped = stripped[end + 3 :].lstrip()
    return stripped


def load_system_prompt(system, system_file):
    if system is not None:
        return system
    if system_file:
        with open(system_file, "r", encoding="utf-8") as f:
            return strip_frontmatter(f.read()).strip()
    raise SystemExit("error: --system か --system-file のいずれかを指定してください")


def load_user_text(input_path):
    if input_path:
        with open(input_path, "r", encoding="utf-8") as f:
            return f.read()
    if sys.stdin.isatty():
        return ""
    # stdin は上で UTF-8 へ reconfigure 済みだが、reconfigure が効かない環境でも
    # 確実に UTF-8 で読めるよう、可能ならバイナリ層から直接復号する(壊れたバイトは置換)。
    buf = getattr(sys.stdin, "buffer", None)
    if buf is not None:
        return buf.read().decode("utf-8", errors="replace")
    return sys.stdin.read()


def load_history(path):
    """会話履歴 JSON を読み、[{"role": "user"|"assistant", "text": str}, ...] に正規化する。
    ステートレス人格(openai/gemini)に過去ラウンドの文脈を多ターンで渡すため(--history)。
    role は user/assistant に丸める(gemini 側で model へ写像)。空テキストは捨てる。"""
    if not path:
        return []
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
    except OSError as e:
        raise SystemExit(f"error: --history を開けません: {e}")
    except ValueError as e:
        raise SystemExit(f"error: --history が不正な JSON です: {e}")
    if not isinstance(data, list):
        raise SystemExit('error: --history は [{"role": ..., "text": ...}, ...] の JSON 配列にしてください')
    norm = []
    for i, m in enumerate(data):
        if not isinstance(m, dict):
            raise SystemExit(f"error: --history[{i}] はオブジェクトではありません")
        role = m.get("role", "user")
        role = "assistant" if role in ("assistant", "model") else "user"
        text = m.get("text", m.get("content", ""))
        if str(text).strip():
            norm.append({"role": role, "text": str(text)})
    return norm


# --------------------------------------------------------------------------- #
# 添付ファイル
# --------------------------------------------------------------------------- #
def classify(path):
    """添付ファイルを 'image' / 'pdf' / 'office' / 'text' / None に分類。"""
    ext = os.path.splitext(path)[1].lower()
    if ext in IMAGE_EXTS:
        return "image"
    if ext in PDF_EXTS:
        return "pdf"
    if ext in OFFICE_EXTS:
        return "office"
    if ext in TEXT_EXTS:
        return "text"
    return None


def guess_mime(path):
    return mimetypes.guess_type(path)[0] or "application/octet-stream"


def file_b64(path):
    with open(path, "rb") as f:
        return base64.b64encode(f.read()).decode("ascii")


def data_url(path):
    mime = guess_mime(path)
    return f"data:{mime};base64,{file_b64(path)}", mime


def read_text_file(path):
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f"=== 添付ファイル: {os.path.basename(path)} ===\n{f.read()}"


def extract_office(path, ext):
    """Office ファイルからテキストを抽出(ネイティブ非対応のための準ネイティブ処理)。"""
    name = os.path.basename(path)
    try:
        if ext == ".docx":
            from docx import Document
            from docx.oxml.ns import qn
            from docx.table import Table
            from docx.text.paragraph import Paragraph

            doc = Document(path)
            lines = []
            # 段落と表を「ドキュメント順」に走査する。doc.paragraphs と doc.tables を別々に
            # 並べると表が本文末尾へまとまり、図表とその文脈の対応が壊れる(資料レビューで不利)。
            for child in doc.element.body.iterchildren():
                if child.tag == qn("w:p"):
                    para = Paragraph(child, doc)
                    if para.text.strip():
                        lines.append(para.text)
                elif child.tag == qn("w:tbl"):
                    tbl = Table(child, doc)
                    for row in tbl.rows:
                        lines.append("\t".join(c.text for c in row.cells))
            body = "\n".join(lines)
        elif ext == ".xlsx":
            import openpyxl

            max_rows = _int_setting("MAGI_XLSX_MAX_ROWS", 2000)
            wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
            chunks, total, truncated = [], 0, False
            for ws in wb.worksheets:
                if truncated:
                    break
                chunks.append(f"# シート: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    if total >= max_rows:  # 暴発防止のソフト上限
                        truncated = True
                        break
                    cells = ["" if v is None else str(v) for v in row]
                    if any(cells):
                        chunks.append("\t".join(cells))
                        total += 1
            if truncated:
                chunks.append(
                    f"…(出力は {max_rows} 行で打ち切り。全体は MAGI_XLSX_MAX_ROWS で調整可)"
                )
            wb.close()  # read_only モードはファイルハンドルを保持する。Windows でのロック回避に明示クローズ
            body = "\n".join(chunks)
        elif ext == ".pptx":
            from pptx import Presentation

            prs = Presentation(path)
            chunks = []
            for i, slide in enumerate(prs.slides, 1):
                chunks.append(f"# スライド {i}")
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        chunks.append(shape.text_frame.text)
            body = "\n".join(chunks)
        else:  # 到達しない
            body = ""
    except ImportError as e:
        pkg = {"docx": "python-docx", "openpyxl": "openpyxl", "pptx": "python-pptx"}.get(
            getattr(e, "name", ""), getattr(e, "name", "")
        )
        raise SystemExit(
            f"error: {ext} の読み取りに {pkg} が必要です。`pip install {pkg}` を実行してください"
        )
    return f"=== 添付ファイル(テキスト抽出): {name} ===\n{body}\n=== ここまで {name} ==="


# --------------------------------------------------------------------------- #
# HTTP ヘルパー(一過性エラーの指数バックオフ・リトライ付き)
# --------------------------------------------------------------------------- #
# 429(レート制限)/ 500・502・503・504(一時的なサーバ障害)や接続タイムアウトは
# 一過性のことが多い。即諦めず指数バックオフ + ジッタで数回リトライする(Retry-After 尊重)。
# 既定は最大4回・各タイムアウト180秒。MAGI_HTTP_MAX_RETRIES / MAGI_HTTP_TIMEOUT で上書き可。
RETRY_STATUSES = frozenset({429, 500, 502, 503, 504})
# フォールバックモデルを試す価値がある HTTP 状態(過負荷=RETRY_STATUSES、モデル不在=404)。
# 401/403(認証)や一般的な 400 は別モデルでも直らないのでフォールバックしない。
FALLBACK_STATUSES = RETRY_STATUSES | frozenset({404})
_BACKOFF_CAP = 30.0


class ProviderHTTPError(Exception):
    """HTTP リクエストが(A2 のリトライ後も)失敗したことを表す。code は HTTP 状態(int)、
    接続エラーのときは None。呼び出し側がフォールバック判断や整形に使う。"""

    def __init__(self, code, detail):
        self.code = code
        self.detail = detail
        super().__init__(f"HTTP {code}: {detail}" if code else f"接続エラー: {detail}")

    @property
    def fallback_worthy(self):
        if not isinstance(self.code, int):
            return False
        if self.code in FALLBACK_STATUSES:
            return True
        # 一部プロバイダ(OpenAI)は不明モデルを 404 でなく 400 + model_not_found で返す。
        # これもモデル切替で回復しうるのでフォールバック対象に含める(一般的な 400 は対象外)。
        return self.code == 400 and ("model_not_found" in self.detail or "does not exist" in self.detail)


def fmt_http_error(provider, err, model=None):
    where = f" (model={model})" if model else ""
    if isinstance(err.code, int):
        return f"error: {provider} HTTP {err.code}{where}: {err.detail}"
    return f"error: {provider} に接続できません{where}: {err.detail}"


def _int_setting(name, default):
    val = get_setting(name)
    try:
        return max(0, int(val))
    except (TypeError, ValueError):
        return default


def _retry_delay(attempt, retry_after):
    """attempt(0 始まり)に対する待機秒。Retry-After(秒)指定があれば優先。"""
    if retry_after is not None:
        try:
            return min(_BACKOFF_CAP, max(0.0, float(retry_after)))
        except (TypeError, ValueError):
            pass
    base = min(_BACKOFF_CAP, 1.0 * (2 ** attempt))
    return base * (0.5 + random.random() * 0.5)  # 50–100% のジッタ


def _backoff(provider, label, attempt, max_retries, retry_after):
    delay = _retry_delay(attempt, retry_after)
    sys.stderr.write(
        f"warn: {provider} {label} — {delay:.1f}s 後に再試行 ({attempt + 1}/{max_retries})\n"
    )
    time.sleep(delay)


def _urlopen_retrying(req, provider):
    """urlopen を実行し、一過性エラーは指数バックオフでリトライして応答を返す。
    リトライ不能・回数超過のエラーは ProviderHTTPError を送出する(呼び出し側で
    フォールバック判断/整形。最終的に未捕捉なら各ブリッジの main が整形して SystemExit)。"""
    max_retries = _int_setting("MAGI_HTTP_MAX_RETRIES", 4)
    timeout = _int_setting("MAGI_HTTP_TIMEOUT", 180)
    attempt = 0
    while True:
        try:
            return urllib.request.urlopen(req, timeout=timeout)
        except urllib.error.HTTPError as e:
            if e.code in RETRY_STATUSES and attempt < max_retries:
                retry_after = e.headers.get("Retry-After") if e.headers else None
                _backoff(provider, f"HTTP {e.code}", attempt, max_retries, retry_after)
                attempt += 1
                continue
            detail = e.read().decode("utf-8", errors="replace")
            raise ProviderHTTPError(e.code, detail)
        except (urllib.error.URLError, TimeoutError, socket.timeout) as e:
            reason = getattr(e, "reason", e)
            transient = isinstance(reason, (TimeoutError, ConnectionError)) or isinstance(
                e, (TimeoutError, socket.timeout)
            )
            if transient and attempt < max_retries:
                _backoff(provider, f"接続エラー({reason})", attempt, max_retries, None)
                attempt += 1
                continue
            raise ProviderHTTPError(None, str(reason))


def http_json(method, url, headers, data_bytes, provider="API"):
    """JSON を返す HTTP 呼び出し。一過性エラーはリトライし、それ以外は整形して SystemExit。"""
    req = urllib.request.Request(url, data=data_bytes, headers=headers, method=method)
    with _urlopen_retrying(req, provider) as resp:
        return json.loads(resp.read().decode("utf-8"))


def http_raw(method, url, headers, data_bytes, provider="API"):
    """生レスポンス(ヘッダ辞書, 本文bytes)を返す。Files API のヘッダ取得用。"""
    req = urllib.request.Request(url, data=data_bytes, headers=headers, method=method)
    with _urlopen_retrying(req, provider) as resp:
        return dict(resp.headers), resp.read()


def run_with_fallback(run_fn, primary, fallback, provider):
    """run_fn(model) を primary で実行。A2 のリトライを使い切ってなお過負荷/モデル不在
    (429/5xx/404)で失敗し、かつ fallback が指定されていれば、代替モデルで1回だけ再試行する。
    人格定義(frontmatter)は変えず、実際に使われたモデルを stderr に明記する(再現性のため)。
    認証エラーや空/拒否応答(A3)はフォールバックの対象外。最終失敗は整形して SystemExit。"""
    try:
        return run_fn(primary)
    except ProviderHTTPError as e:
        if fallback and fallback != primary and e.fallback_worthy:
            sys.stderr.write(
                f"warn: {provider} の model={primary} が失敗(HTTP {e.code})。"
                f"fallback model={fallback} で再試行します"
                f"(成果物には実際に使われたモデル={fallback} を明記すること)\n"
            )
            try:
                return run_fn(fallback)
            except ProviderHTTPError as e2:
                raise SystemExit(fmt_http_error(provider, e2, fallback))
        raise SystemExit(fmt_http_error(provider, e, primary))


# --------------------------------------------------------------------------- #
# 生成レベルのリトライ(空応答・拒否応答)
# --------------------------------------------------------------------------- #
# HTTP は 200 でも、モデルが空文字や拒否(セーフティ)を返すことがある。これは A2 の
# HTTP リトライとは別レイヤなので、call 全体を MAGI_GEN_MAX_RETRIES 回まで再試行する。
class EmptyOrRefusalResponse(Exception):
    """モデルが空応答または拒否応答を返したことを示す(再試行候補)。
    kind は 'empty' か 'refusal'、detail は原因の手掛かり。"""

    def __init__(self, kind, detail=""):
        self.kind = kind
        self.detail = detail
        super().__init__(f"{kind}: {detail}")


def run_with_retry(call_fn, provider):
    """call_fn()(API 呼び出し1回ぶん)を実行し、空/拒否応答なら同じ要求で再試行する。
    既定は1回(計2回試行)。MAGI_GEN_MAX_RETRIES で調整可。最終的に失敗したら SystemExit。
    一過性の HTTP エラーは call_fn 内(http_json)の A2 リトライで別途吸収される。"""
    max_retries = _int_setting("MAGI_GEN_MAX_RETRIES", 1)
    attempt = 0
    while True:
        try:
            return call_fn()
        except EmptyOrRefusalResponse as e:
            label = "空応答" if e.kind == "empty" else "拒否応答"
            if attempt < max_retries:
                sys.stderr.write(
                    f"warn: {provider} が{label}を返しました — 同じ要求で再試行 "
                    f"({attempt + 1}/{max_retries})\n"
                )
                attempt += 1
                continue
            raise SystemExit(
                f"error: {provider} の{label}が続きました。人格定義は変えず、"
                f"指示(プロンプト)を言い換えて再実行してください: {e.detail}"
            )


def encode_multipart(fields, file_field, filename, file_bytes, content_type):
    """multipart/form-data を手組み(stdlib だけで Files API にアップロードするため)。"""
    boundary = "----magiboundary7MA4YWxkTrZu0gW"
    crlf = b"\r\n"
    body = []
    for key, value in fields.items():
        body.append(b"--" + boundary.encode())
        body.append(f'Content-Disposition: form-data; name="{key}"'.encode())
        body.append(b"")
        body.append(str(value).encode())
    body.append(b"--" + boundary.encode())
    body.append(
        f'Content-Disposition: form-data; name="{file_field}"; filename="{filename}"'.encode()
    )
    body.append(f"Content-Type: {content_type}".encode())
    body.append(b"")
    body.append(file_bytes)
    body.append(b"--" + boundary.encode() + b"--")
    body.append(b"")
    return crlf.join(body), f"multipart/form-data; boundary={boundary}"
