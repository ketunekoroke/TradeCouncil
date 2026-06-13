#!/usr/bin/env python3
"""extract_office.py — Office ファイル(docx/pptx/xlsx)を本文順の Markdown に抽出する CLI。

資料レビュー(document-review)で、対象資料を全人格に等しく渡す「正本テキスト」を作るための
ヘルパー。bridge_common の Office 抽出が LLM 注入用のフラットなテキストなのに対し、こちらは
見出し(#/##/###)・GFM 表・画像マーカーを保った Markdown を出力し、人が読める/改訂の土台に
できる形にする。docx は段落と表をドキュメント順に並べる(→ A5 と同じ原則)。

使い方:
  python scripts/extract_office.py 資料.docx                # stdout に Markdown
  python scripts/extract_office.py 資料.docx -o canon.md    # ファイルへ
  python scripts/extract_office.py スライド.pptx
  python scripts/extract_office.py 表.xlsx

依存: python-docx(docx)/ python-pptx(pptx)/ openpyxl(xlsx)。必要時のみ。
環境変数: MAGI_XLSX_MAX_ROWS(xlsx の出力行ソフト上限。既定 2000)。
"""
import argparse
import os
import sys

for _s in (sys.stdout, sys.stderr):
    try:
        _s.reconfigure(encoding="utf-8")
    except (AttributeError, ValueError):
        pass


def _need(pkg, ext):
    raise SystemExit(f"error: {ext} の読み取りに {pkg} が必要です。`pip install {pkg}` を実行してください")


def _cell(c):
    return c.text.replace("\n", " / ").strip()


def docx_to_md(path):
    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError:
        _need("python-docx", ".docx")
    doc = Document(path)
    lines = []
    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            para = Paragraph(child, doc)
            style = para.style.name if para.style else ""
            text = para.text.strip()
            if para._p.findall(".//" + qn("w:drawing")):
                lines.append("*[画像]*")
            if not text:
                continue
            if style == "Heading 1":
                lines.append("\n# " + text)
            elif style == "Heading 2":
                lines.append("\n## " + text)
            elif style == "Heading 3":
                lines.append("\n### " + text)
            else:
                lines.append(text)
        elif child.tag == qn("w:tbl"):
            tbl = Table(child, doc)
            rows = tbl.rows
            if not rows:
                continue
            header = [_cell(c) for c in rows[0].cells]
            lines.append("")
            lines.append("| " + " | ".join(header) + " |")
            lines.append("| " + " | ".join(["---"] * len(header)) + " |")
            for r in rows[1:]:
                lines.append("| " + " | ".join(_cell(c) for c in r.cells) + " |")
            lines.append("")
    return "\n".join(lines).strip() + "\n"


def pptx_to_md(path):
    try:
        from pptx import Presentation
    except ImportError:
        _need("python-pptx", ".pptx")
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n# スライド {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
    return "\n".join(lines).strip() + "\n"


def xlsx_to_md(path):
    try:
        import openpyxl
    except ImportError:
        _need("openpyxl", ".xlsx")
    try:
        max_rows = int(os.environ.get("MAGI_XLSX_MAX_ROWS", "2000"))
    except (TypeError, ValueError):
        max_rows = 2000
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    lines, total = [], 0
    for ws in wb.worksheets:
        lines.append(f"\n## シート: {ws.title}")
        printed_header = False
        for row in ws.iter_rows(values_only=True):
            if total >= max_rows:
                lines.append(f"…(出力は {max_rows} 行で打ち切り。全体は MAGI_XLSX_MAX_ROWS で調整可)")
                break
            cells = ["" if v is None else str(v) for v in row]
            if not any(cells):
                continue
            lines.append("| " + " | ".join(cells) + " |")
            if not printed_header:
                lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
                printed_header = True
            total += 1
    return "\n".join(lines).strip() + "\n"


EXTRACTORS = {".docx": docx_to_md, ".pptx": pptx_to_md, ".xlsx": xlsx_to_md}


def main():
    p = argparse.ArgumentParser(description="Office(docx/pptx/xlsx)を本文順の Markdown に抽出")
    p.add_argument("path", help="入力ファイル(.docx/.pptx/.xlsx)")
    p.add_argument("-o", "--output", help="出力 Markdown のパス(省略時は stdout)")
    a = p.parse_args()
    ext = os.path.splitext(a.path)[1].lower()
    fn = EXTRACTORS.get(ext)
    if fn is None:
        raise SystemExit(f"error: 未対応の形式です: {ext}(対応: .docx / .pptx / .xlsx)")
    if not os.path.exists(a.path):
        raise SystemExit(f"error: ファイルが見つかりません: {a.path}")
    md = fn(a.path)
    if a.output:
        with open(a.output, "w", encoding="utf-8") as f:
            f.write(md)
        sys.stderr.write(f"wrote: {a.output} ({len(md)} chars)\n")
    else:
        sys.stdout.write(md)


if __name__ == "__main__":
    main()
